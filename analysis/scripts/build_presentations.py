"""
Build both PhD presentations from scratch using python-pptx.

Department (Semester 2 progress):
  Title → Intro → Hypotheses → Objectives →
  Sem 1 Progress → Sem 2 Progress → Results (4 figs) → Next Steps

Joint Seminar:
  Title → Intro → Objectives → Hypotheses →
  Results (6 figs) → Key Findings → Next Steps
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ─────────────────────────────────────────────────────────────────────
FIG   = '/home/abdul/TESTs/WRFGC_SEAsia_Jan2025/analysis/figures'
OUT   = '/home/abdul/TESTs/WRFGC_SEAsia_Jan2025/analysis/presentations'

FIGS = {
    'four_panel':  os.path.join(FIG, 'fig01_fourpanel_species.png'),
    'pblh_vc':     os.path.join(FIG, 'fig02_pblh_vc.png'),
    'stag':        os.path.join(FIG, 'fig03_stagnation_composite.png'),
    'timeseries':  os.path.join(FIG, 'fig04_timeseries.png'),
    'o3_wind':     os.path.join(FIG, 'fig05_o3_wind.png'),
    'domain':      os.path.join(FIG, 'fig06_domain_overview.png'),
}

# ── Colours / Fonts ───────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x34, 0x5C)
TEAL   = RGBColor(0x00, 0x7A, 0x8A)
ORANGE = RGBColor(0xD4, 0x6B, 0x08)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF2, 0xF4, 0xF8)
DGREY  = RGBColor(0x44, 0x44, 0x44)
BLACK  = RGBColor(0x11, 0x11, 0x11)


# ══════════════════════════════════════════════════════════════════════════════
# LOW-LEVEL HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def new_prs(wide=True):
    """Create a new 16:9 presentation."""
    prs = Presentation()
    if wide:
        prs.slide_width  = Emu(9_144_000)   # 10 inches
        prs.slide_height = Emu(5_143_500)   # 5.625 inches
    return prs


def blank_slide(prs, bg_color=None):
    """Add a blank slide; optionally fill background."""
    layout = prs.slide_layouts[6]   # truly blank
    slide  = prs.slides.add_slide(layout)
    if bg_color:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return slide


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    """Add a filled rectangle shape."""
    from pptx.util import Pt as PtU
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = PtU(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def txb(slide, left, top, width, height):
    """Add a text-box and return (shape, text_frame)."""
    shape = slide.shapes.add_textbox(
        Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    tf = shape.text_frame
    tf.word_wrap = True
    return shape, tf


def para(tf, text, size_pt=18, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, space_before=0, space_after=0,
         level=0, italic=False):
    """Add a paragraph to a text frame."""
    from pptx.util import Pt as PtU
    p = tf.add_paragraph()
    p.text  = text
    p.level = level
    p.alignment = align
    p.space_before = PtU(space_before)
    p.space_after  = PtU(space_after)
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size   = PtU(size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def first_para(tf, text, size_pt=18, bold=False, color=BLACK,
               align=PP_ALIGN.LEFT, italic=False):
    """Set the first (auto-created) paragraph of a text frame."""
    from pptx.util import Pt as PtU
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size   = PtU(size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_img(slide, path, left, top, width, height):
    """Add image to slide."""
    slide.shapes.add_picture(
        path,
        Emu(int(left)), Emu(int(top)),
        width=Emu(int(width)), height=Emu(int(height)),
    )


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE TEMPLATE HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def W(prs): return prs.slide_width
def H(prs): return prs.slide_height


def add_header_bar(slide, prs, title_text, subtitle_text=''):
    """Dark navy bar at top with white title + optional subtitle."""
    bar_h = H(prs) * 0.16
    add_rect(slide, 0, 0, W(prs), bar_h, fill_rgb=NAVY)

    # Title
    _, tf = txb(slide, W(prs)*0.03, H(prs)*0.01, W(prs)*0.94, bar_h*0.6)
    first_para(tf, title_text, size_pt=22, bold=True, color=WHITE)

    # Subtitle
    if subtitle_text:
        _, tf2 = txb(slide, W(prs)*0.03, H(prs)*0.01 + bar_h*0.55,
                     W(prs)*0.94, bar_h*0.4)
        first_para(tf2, subtitle_text, size_pt=13, color=RGBColor(0xC8, 0xD8, 0xF0))


def add_footer(slide, prs, text='Abdul | PhD Semester 2 | March 2026'):
    """Thin footer bar."""
    bar_h = H(prs) * 0.055
    add_rect(slide, 0, H(prs) - bar_h, W(prs), bar_h,
             fill_rgb=RGBColor(0xE8, 0xEC, 0xF4))
    _, tf = txb(slide, W(prs)*0.02, H(prs) - bar_h + bar_h*0.1,
                W(prs)*0.96, bar_h*0.8)
    first_para(tf, text, size_pt=9, color=RGBColor(0x66, 0x66, 0x66),
               align=PP_ALIGN.CENTER)


def content_slide(prs, title, subtitle='', footer=True, footer_text='Abdul | PhD Semester 2 | March 2026'):
    """Create a standard content slide with header bar + footer. Returns slide."""
    slide = blank_slide(prs, bg_color=WHITE)
    add_header_bar(slide, prs, title, subtitle)
    if footer:
        add_footer(slide, prs, footer_text)
    return slide


def figure_slide(prs, img_path, title, subtitle='',
                 img_top_frac=0.17, img_height_frac=0.76,
                 footer_text='Abdul | PhD Semester 2 | March 2026'):
    """Full-width figure slide."""
    slide = content_slide(prs, title, subtitle, footer_text=footer_text)
    img_top  = H(prs) * img_top_frac
    img_left = W(prs) * 0.02
    img_w    = W(prs) * 0.96
    img_h    = H(prs) * img_height_frac
    add_img(slide, img_path, img_left, img_top, img_w, img_h)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# BULLET HELPER
# ══════════════════════════════════════════════════════════════════════════════

def bullet_block(slide, prs, items, left_frac=0.05, top_frac=0.20,
                 width_frac=0.90, height_frac=0.68,
                 main_size=16, sub_size=14, main_color=NAVY, sub_color=DGREY):
    """
    items: list of (text, level) tuples.
      level 0 → main bullet, level 1 → sub-bullet
    """
    _, tf = txb(slide,
                W(prs)*left_frac, H(prs)*top_frac,
                W(prs)*width_frac, H(prs)*height_frac)
    first = True
    for text, level in items:
        size  = main_size if level == 0 else sub_size
        color = main_color if level == 0 else sub_color
        bold  = (level == 0)
        sp_b  = 6 if level == 0 else 2
        if first:
            first_para(tf, text, size_pt=size, bold=bold, color=color)
            first = False
        else:
            p = tf.add_paragraph()
            p.text  = text
            p.level = level
            p.space_before = Pt(sp_b)
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# ① DEPARTMENT PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════

def build_department():
    prs = new_prs()
    ft  = 'Abdul | PhD Sem 2 Progress | March 2026'

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = blank_slide(prs, bg_color=NAVY)

    # Decorative teal strip on left
    add_rect(slide, 0, 0, W(prs)*0.008, H(prs), fill_rgb=TEAL)

    _, tf = txb(slide, W(prs)*0.07, H(prs)*0.22, W(prs)*0.86, H(prs)*0.28)
    first_para(tf,
        'Semester 2 PhD Progress Seminar',
        size_pt=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    _, tf2 = txb(slide, W(prs)*0.07, H(prs)*0.50, W(prs)*0.86, H(prs)*0.16)
    first_para(tf2,
        'Meteorology–Chemistry Coupling and PM₂.₅/SOA Variability over Thailand',
        size_pt=17, bold=False, color=RGBColor(0xC8, 0xD8, 0xF0),
        align=PP_ALIGN.LEFT)

    _, tf3 = txb(slide, W(prs)*0.07, H(prs)*0.68, W(prs)*0.86, H(prs)*0.20)
    for i, line in enumerate([
        'Abdul  |  Sustainable & Environmental Engineering',
        'Supervisor: [Supervisor Name]',
        'March 2026',
    ]):
        if i == 0:
            first_para(tf3, line, size_pt=13, color=RGBColor(0xAA, 0xCC, 0xEE))
        else:
            para(tf3, line, size_pt=12, color=RGBColor(0x88, 0xAA, 0xCC))

    # ── Slide 2: Introduction ─────────────────────────────────────────────────
    slide = content_slide(prs, 'Introduction', 'Context and motivation', footer_text=ft)
    bullet_block(slide, prs, [
        ('Research context', 0),
        ('PM₂.₅ is a leading cause of premature mortality in Southeast Asia', 1),
        ('Thailand experiences severe haze episodes driven by biomass burning, traffic, and stagnant meteorology', 1),
        ('Offline chemistry models (GEOS-Chem) ignore two-way meteorology–chemistry feedbacks', 1),
        ('', 0),
        ('Research gap', 0),
        ('No study has applied WRF-GC (online coupled WRF + GEOS-Chem) over Thailand', 1),
        ('The role of stagnation episodes in modulating SOA formation is poorly quantified', 1),
        ('NO₃ radical pathways under Bangkok high-NOx nights remain uncharacterised', 1),
        ('', 0),
        ('This work', 0),
        ('First WRF-GC simulation over South/Southeast Asia — January 2025', 1),
        ('27 km domain (179×149), online two-way coupled chemistry', 1),
    ], main_size=17, sub_size=14)

    # ── Slide 3: Hypotheses ───────────────────────────────────────────────────
    slide = content_slide(prs, 'Research Hypotheses', footer_text=ft)

    boxes = [
        ('H1', 'WRF-GC reduces PM₂.₅ RMSE/NMB by 20–40% compared to offline GEOS-Chem over Thailand',
         TEAL),
        ('H2', 'Local anthropogenic SOA fraction increases by ≥15–20% under stagnant conditions '
               '(VC < 4000 m² s⁻¹) due to prolonged precursor residence time',
         ORANGE),
        ('H3', 'Stagnant high-NOx conditions shift the dominant SOA initiation pathway '
               'from OH-initiated to NO₃-initiated oxidation',
         RGBColor(0x6A, 0x0D, 0x83)),
    ]

    top_frac  = 0.20
    box_h     = H(prs) * 0.20
    gap       = H(prs) * 0.03

    for b_label, b_text, b_color in boxes:
        # Colour tab on left
        add_rect(slide, W(prs)*0.05, top_frac,
                 W(prs)*0.012, box_h, fill_rgb=b_color)
        # Light bg
        add_rect(slide, W(prs)*0.062, top_frac,
                 W(prs)*0.88, box_h,
                 fill_rgb=RGBColor(0xF5, 0xF7, 0xFC),
                 line_rgb=RGBColor(0xCC, 0xD5, 0xE8), line_width_pt=0.5)
        # Label
        _, tfl = txb(slide, W(prs)*0.065, top_frac + box_h*0.08,
                     W(prs)*0.06, box_h*0.5)
        first_para(tfl, b_label, size_pt=18, bold=True, color=b_color)
        # Text
        _, tft = txb(slide, W(prs)*0.13, top_frac + box_h*0.08,
                     W(prs)*0.80, box_h*0.84)
        first_para(tft, b_text, size_pt=14, color=DGREY)
        top_frac += box_h + gap

    add_footer(slide, prs, ft)

    # ── Slide 4: Objectives ───────────────────────────────────────────────────
    slide = content_slide(prs, 'Research Objectives', footer_text=ft)
    bullet_block(slide, prs, [
        ('Objective 1 — Model Performance Evaluation', 0),
        ('Compare WRF-GC vs offline GEOS-Chem PM₂.₅ against observational data', 1),
        ('Quantify improvement in RMSE and NMB (target: 20–40% reduction)', 1),
        ('', 0),
        ('Objective 2 — Regime-Dependent SOA Attribution', 0),
        ('Attribute SOA sources under stagnant (VC < 4000 m² s⁻¹) vs ventilated regimes', 1),
        ('Quantify changes in local anthropogenic vs biogenic SOA fractions', 1),
        ('', 0),
        ('Objective 3 — Nighttime NO₃ Chemistry Pathways', 0),
        ('Characterise NO₃ radical chemistry in Bangkok urban high-NOx environment', 1),
        ('Identify dominant SOA initiation pathways during stagnant nights', 1),
    ], main_size=17, sub_size=14)

    # ── Slide 5: Semester 1 Progress ─────────────────────────────────────────
    slide = content_slide(prs, 'Semester 1 Progress',
                          'Domain configuration, data acquisition, model setup', footer_text=ft)
    bullet_block(slide, prs, [
        ('Domain Design & Configuration', 0),
        ('Designed d01 outer domain: 27 km, 179×149 grid covering South & Southeast Asia', 1),
        ('Configured nested d02 domain over Thailand (9 km) for high-resolution analysis', 1),
        ('Set up WRF physics options: WSM6 microphysics, YSU PBL, RRTMG radiation', 1),
        ('', 0),
        ('Meteorological Data', 0),
        ('Downloaded and processed ERA5 reanalysis data (0.25° × 0.25°, Jan 2025)', 1),
        ('Ran WPS preprocessing: geogrid, ungrib, metgrid for ICBC preparation', 1),
        ('', 0),
        ('Model Installation & Compilation', 0),
        ('Compiled WRF-GC v2.0 on HPC cluster (Intel MPI, HDF5/NetCDF libraries)', 1),
        ('Resolved library dependency issues and validated test runs', 1),
        ('Performed initial short test simulations to verify model stability', 1),
    ], main_size=16, sub_size=13)

    # ── Slide 6: Semester 2 Progress ─────────────────────────────────────────
    slide = content_slide(prs, 'Semester 2 Progress',
                          'Hypothesis development, actual simulation, proposal preparation', footer_text=ft)
    bullet_block(slide, prs, [
        ('Hypothesis Development', 0),
        ('Developed three testable hypotheses linking meteorology–chemistry coupling to PM₂.₅/SOA', 1),
        ('Conducted literature review — identified WRF-GC advantage over offline models', 1),
        ('Refined research questions based on Thailand-specific air quality challenges', 1),
        ('', 0),
        ('Actual WRF-GC Simulation', 0),
        ('Completed full production run: Jan 3–28, 2025 (26-day simulation)', 1),
        ('Validated met output (temperature, winds, PBLH) against ERA5 climatology', 1),
        ('Extracted surface species: O₃, NO₂, CO, PM₂.₅, PBLH, ventilation coefficient', 1),
        ('Identified stagnation events (VC < 4000 m² s⁻¹) — dominated January 2025', 1),
        ('', 0),
        ('Proposal Exam Preparation', 0),
        ('Drafted full PhD proposal with objectives, hypotheses, and methodology', 1),
        ('Preparing for proposal examination (target: end of Semester 2)', 1),
    ], main_size=16, sub_size=13)

    # ── Slides 7–10: Results ──────────────────────────────────────────────────
    figure_slide(prs, FIGS['four_panel'],
                 'WRF-GC Results — Surface Species (January 2025 Mean)',
                 'O₃, NO₂, CO, PM₂.₅ at surface level | 27 km domain',
                 footer_text=ft)

    figure_slide(prs, FIGS['stag'],
                 'PM₂.₅ Stagnation Composite Analysis',
                 'Strong stagnation (VC < 1243 m² s⁻¹) vs weak stagnation — January 2025',
                 footer_text=ft)

    figure_slide(prs, FIGS['timeseries'],
                 'Daily Time Series — Bangkok & Thailand Mean Concentrations',
                 'O₃, NO₂, CO, PM₂.₅ with WHO 2021 and Thailand NAAQS reference lines',
                 footer_text=ft)

    figure_slide(prs, FIGS['pblh_vc'],
                 'Mixing Diagnostics — PBLH and Ventilation Coefficient',
                 'January 2025 mean | Low VC over northern Thailand indicates stagnation-prone areas',
                 footer_text=ft)

    # ── Slide 11: Next Steps ──────────────────────────────────────────────────
    slide = content_slide(prs, 'Next Steps & Timeline', footer_text=ft)
    bullet_block(slide, prs, [
        ('Immediate (Next 1–2 months)', 0),
        ('Collect observational PM₂.₅ data for model–obs comparison (Obj 1)', 1),
        ('Run offline GEOS-Chem simulation for same period as WRF-GC baseline', 1),
        ('Submit PhD proposal for examination', 1),
        ('', 0),
        ('Medium Term (Semester 3)', 0),
        ('Perform SOA source apportionment analysis using GEOS-Chem tracers (Obj 2)', 1),
        ('Analyse NO₃ radical production/loss budgets for Bangkok (Obj 3)', 1),
        ('Extend simulation to full dry season (Oct 2024 – March 2025)', 1),
        ('', 0),
        ('Longer Term', 0),
        ('Write and submit first manuscript: WRF-GC PM₂.₅ evaluation over Thailand', 1),
        ('Conference presentation of results', 1),
    ], main_size=16, sub_size=13)

    out_path = os.path.join(OUT, 'PhD_department_Sem2_v2.pptx')
    prs.save(out_path)
    print(f'Department pptx saved: {out_path}  ({len(prs.slides)} slides)')
    return out_path


# ══════════════════════════════════════════════════════════════════════════════
# ② JOINT SEMINAR PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════

def build_joint():
    prs = new_prs()
    ft  = 'Abdul | WRF-GC Thailand | Joint Seminar 2026'

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = blank_slide(prs, bg_color=RGBColor(0x0D, 0x1B, 0x3E))

    # Teal diagonal-ish accent bar
    add_rect(slide, 0, H(prs)*0.72, W(prs), H(prs)*0.006, fill_rgb=TEAL)

    _, tf = txb(slide, W(prs)*0.06, H(prs)*0.12, W(prs)*0.88, H(prs)*0.30)
    first_para(tf,
        'Role of Meteorology–Chemistry Coupling in Regulating\n'
        'PM₂.₅ and SOA Variability over Thailand',
        size_pt=24, bold=True, color=WHITE)

    _, tf2 = txb(slide, W(prs)*0.06, H(prs)*0.46, W(prs)*0.88, H(prs)*0.12)
    first_para(tf2,
        'First WRF-GC (online coupled) simulation results — January 2025',
        size_pt=15, color=RGBColor(0xC0, 0xD8, 0xF0))

    _, tf3 = txb(slide, W(prs)*0.06, H(prs)*0.60, W(prs)*0.88, H(prs)*0.20)
    for i, line in enumerate([
        'Abdul  |  Sustainable & Environmental Engineering',
        'Joint Seminar  |  March 2026',
    ]):
        if i == 0:
            first_para(tf3, line, size_pt=13, color=RGBColor(0xAA, 0xCC, 0xEE))
        else:
            para(tf3, line, size_pt=12, color=RGBColor(0x88, 0xAA, 0xCC))

    # ── Slide 2: Introduction ─────────────────────────────────────────────────
    slide = content_slide(prs, 'Introduction', footer_text=ft)
    bullet_block(slide, prs, [
        ('Problem Statement', 0),
        ('PM₂.₅ pollution causes ~80,000 premature deaths/year in Southeast Asia', 1),
        ('Thailand regularly exceeds WHO PM₂.₅ 24-hr guideline (15 µg/m³)', 1),
        ('Biomass burning, traffic, and industrial emissions combined with', 1),
        ('stagnant winter meteorology drive severe haze episodes', 1),
        ('', 0),
        ('The modelling challenge', 0),
        ('Offline chemistry models (GEOS-Chem standalone) treat meteorology as fixed', 1),
        ('Cannot capture aerosol–radiation–cloud feedbacks or PBL height changes', 1),
        ('WRF-GC couples WRF (NWP) with GEOS-Chem (chemistry) — two-way online', 1),
        ('', 0),
        ('Study period', 0),
        ('January 2025 — predominantly stagnant; ideal test case for stagnation–PM₂.₅ linkage', 1),
    ], main_size=16, sub_size=13)

    # ── Slide 3: Objectives ───────────────────────────────────────────────────
    slide = content_slide(prs, 'Research Objectives', footer_text=ft)

    objs = [
        ('Obj 1', 'Model Performance Evaluation',
         'Quantify WRF-GC improvement over offline GEOS-Chem for PM₂.₅ (RMSE, NMB, correlation)',
         TEAL),
        ('Obj 2', 'Regime-Dependent SOA Attribution',
         'Attribute SOA sources under stagnant vs ventilated regimes; quantify local anthropogenic SOA change',
         ORANGE),
        ('Obj 3', 'Nighttime NO₃ Chemistry',
         'Characterise NO₃ radical pathways in Bangkok high-NOx environment and identify dominant SOA initiation routes',
         RGBColor(0x6A, 0x0D, 0x83)),
    ]

    top_frac = 0.19
    box_h    = H(prs) * 0.215
    gap      = H(prs) * 0.025

    for o_label, o_title, o_text, o_color in objs:
        add_rect(slide, W(prs)*0.04, top_frac, W(prs)*0.01, box_h, fill_rgb=o_color)
        add_rect(slide, W(prs)*0.05, top_frac, W(prs)*0.90, box_h,
                 fill_rgb=RGBColor(0xF5, 0xF7, 0xFC),
                 line_rgb=RGBColor(0xCC, 0xD5, 0xE8), line_width_pt=0.5)
        _, tfl = txb(slide, W(prs)*0.055, top_frac + box_h*0.06, W(prs)*0.07, box_h*0.45)
        first_para(tfl, o_label, size_pt=15, bold=True, color=o_color)
        _, tft = txb(slide, W(prs)*0.13, top_frac + box_h*0.06, W(prs)*0.80, box_h*0.9)
        first_para(tft, o_title, size_pt=14, bold=True, color=NAVY)
        para(tft, o_text, size_pt=13, color=DGREY, space_before=4)
        top_frac += box_h + gap

    add_footer(slide, prs, ft)

    # ── Slide 4: Hypotheses ───────────────────────────────────────────────────
    slide = content_slide(prs, 'Hypotheses', footer_text=ft)
    bullet_block(slide, prs, [
        ('H1 — Model Performance', 0),
        ('WRF-GC reduces PM₂.₅ RMSE/NMB by 20–40% compared to offline GEOS-Chem', 1),
        ('Rationale: Online coupling captures PBL feedbacks that alter vertical mixing of pollutants', 1),
        ('', 0),
        ('H2 — SOA Stagnation Response', 0),
        ('Local anthropogenic SOA fraction increases by ≥15–20% under stagnant conditions', 1),
        ('(Ventilation Coefficient < 4000 m² s⁻¹)', 1),
        ('Rationale: Longer precursor residence times enhance secondary aerosol formation', 1),
        ('', 0),
        ('H3 — Nighttime Chemistry Shift', 0),
        ('Stagnant high-NOx conditions shift dominant SOA pathway: OH → NO₃-initiated oxidation', 1),
        ('Rationale: Elevated NO₂ increases NO₃ production; NO₃ is a major nighttime oxidant for VOCs', 1),
    ], main_size=16, sub_size=13)

    # ── Slide 5: WRF-GC Setup ────────────────────────────────────────────────
    slide = content_slide(prs, 'WRF-GC Model Setup', 'Online coupled meteorology–chemistry', footer_text=ft)
    bullet_block(slide, prs, [
        ('Model', 0),
        ('WRF-GC v2.0 — WRF 4.x coupled with GEOS-Chem 12.9 (two-way online)', 1),
        ('', 0),
        ('Domain', 0),
        ('d01: 27 km resolution, 179×149 grid points', 1),
        ('Coverage: South & Southeast Asia (lat −4.6–30.7°N, lon 78.3–122.7°E)', 1),
        ('Simulation period: January 3–28, 2025 (7-day spin-up)', 1),
        ('', 0),
        ('Physics & Chemistry', 0),
        ('Microphysics: WSM6  |  PBL: YSU  |  Radiation: RRTMG', 1),
        ('Chemistry: GEOS-Chem full tropospheric + aerosol (ISORROPIA II)', 1),
        ('Emissions: CEDS anthropogenic + GFAS biomass burning + MEGAN biogenic', 1),
        ('', 0),
        ('Boundary conditions', 0),
        ('Meteorology: ERA5 reanalysis (6-hourly, 0.25°)', 1),
        ('Chemistry ICs/BCs: GEOS-Chem 0.25° × 0.3125° global run', 1),
    ], main_size=15, sub_size=13)

    # ── Slides 6–11: Results ──────────────────────────────────────────────────
    figure_slide(prs, FIGS['domain'],
                 'Domain Overview — CO, PBLH, and Ventilation Coefficient',
                 'January 2025 mean | 27 km, South & Southeast Asia',
                 footer_text=ft)

    figure_slide(prs, FIGS['four_panel'],
                 'Surface Species Concentrations — January 2025 Mean',
                 'O₃, NO₂, CO, PM₂.₅ | WRF-GC simulation output',
                 footer_text=ft)

    figure_slide(prs, FIGS['pblh_vc'],
                 'Planetary Boundary Layer Height and Ventilation Coefficient',
                 'Low VC over northern Thailand and Indochina — stagnation-prone region',
                 footer_text=ft)

    figure_slide(prs, FIGS['stag'],
                 'PM₂.₅ Response to Stagnation — Composite Analysis',
                 'Strong stagnation (VC < 1243 m² s⁻¹) vs weak stagnation — January 2025',
                 footer_text=ft)

    figure_slide(prs, FIGS['timeseries'],
                 'Bangkok & Thailand-Mean Daily Concentrations',
                 'With WHO 2021 and Thailand NAAQS guideline thresholds',
                 footer_text=ft)

    figure_slide(prs, FIGS['o3_wind'],
                 'Surface O₃ and 10 m Wind Field',
                 'O₃ elevated over marine areas; southwesterly flow over Bay of Bengal',
                 footer_text=ft)

    # ── Slide 12: Key Findings ───────────────────────────────────────────────
    slide = content_slide(prs, 'Key Findings — Preliminary', footer_text=ft)
    bullet_block(slide, prs, [
        ('Spatial patterns', 0),
        ('PM₂.₅ elevated in northern Thailand and Indochina — consistent with biomass burning transport', 1),
        ('O₃ enhanced over marine/remote areas due to reduced NOx titration', 1),
        ('PBLH lowest over mountainous northern Thailand — limits vertical dilution', 1),
        ('', 0),
        ('Stagnation dominates January 2025', 0),
        ('All 19 post-spin-up days fall below VC = 4000 m² s⁻¹ threshold', 1),
        ('Strong stagnation days (VC < 1243 m² s⁻¹) show ~40–60% higher PM₂.₅ than weak stagnation', 1),
        ('Supports H2: stagnation accumulates aerosol precursors', 1),
        ('', 0),
        ('Bangkok time series', 0),
        ('PM₂.₅ reaches ~150–190 µg/m³ peak at Bangkok grid point in late January', 1),
        ('Both Bangkok and Thailand-mean exceed WHO 24-hr guideline (15 µg/m³)', 1),
        ('Thailand 24-hr NAAQS (37.5 µg/m³) exceeded on multiple days', 1),
    ], main_size=15, sub_size=13)

    # ── Slide 13: Next Steps ─────────────────────────────────────────────────
    slide = content_slide(prs, 'Next Steps', footer_text=ft)
    bullet_block(slide, prs, [
        ('Model evaluation (Obj 1)', 0),
        ('Collect Thai EPA / AERONET / ground station PM₂.₅ observations', 1),
        ('Run offline GEOS-Chem baseline for head-to-head comparison', 1),
        ('', 0),
        ('SOA attribution (Obj 2)', 0),
        ('Use GEOS-Chem tagged tracers to separate biogenic vs anthropogenic SOA', 1),
        ('Compare stagnant vs ventilated day composites (extend to full dry season)', 1),
        ('', 0),
        ('NO₃ chemistry (Obj 3)', 0),
        ('Extract NO₃ production/loss rates and nighttime VOC oxidation budgets', 1),
        ('Focus on Bangkok urban grid — high-NOx / high-VOC urban environment', 1),
        ('', 0),
        ('Publications & Milestones', 0),
        ('PhD proposal examination — Semester 2 (target: May 2026)', 1),
        ('First manuscript submission — Semester 3 (WRF-GC evaluation paper)', 1),
    ], main_size=15, sub_size=13)

    out_path = os.path.join(OUT, 'PhD_joint_seminar_v2.pptx')
    prs.save(out_path)
    print(f'Joint seminar pptx saved: {out_path}  ({len(prs.slides)} slides)')
    return out_path


# ══════════════════════════════════════════════════════════════════════════════
if __name__ == '__main__':
    build_department()
    build_joint()
    print('\nDone.')
